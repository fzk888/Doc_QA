#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
文件处理和分词工具模块
支持多种文件格式的解析和分词处理
"""

import os
import re
from typing import List, Tuple
from langchain.schema import Document

# 尝试导入各种文件处理库
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from add.morefile.rag.nlp.rag_tokenizer import tokenizer
    tokenize = tokenizer.jieba_tokenize
except ImportError:
    # 如果无法导入自定义分词器，使用简单正则表达式分词
    def tokenize(text):
        return " ".join(re.findall(r'[\w]+', text))


class FileProcessor:
    """
    文件处理器类，支持多种文件格式的解析和分词
    """

    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xlsx,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.txt': self._process_txt,
            '.md': self._process_txt,
        }

    def process_file(self, file_path: str) -> List[Document]:
        """
        处理文件并返回分词后的文档列表
        
        Args:
            file_path: 文件路径
            
        Returns:
            List[Document]: 处理后的文档列表
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        _, ext = os.path.splitext(file_path.lower())
        
        if ext in self.supported_formats:
            return self.supported_formats[ext](file_path)
        else:
            # 默认按文本文件处理
            return self._process_txt(file_path)

    def _process_pdf(self, file_path: str) -> List[Document]:
        """处理PDF文件"""
        if PyPDF2 is None:
            raise ImportError("缺少PyPDF2库，无法处理PDF文件")
            
        documents = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        # 对文本进行分词
                        tokenized_text = tokenize(text)
                        documents.append(Document(
                            page_content=tokenized_text,
                            metadata={
                                'file_path': file_path,
                                'page': page_num + 1,
                                'file_type': 'pdf'
                            }
                        ))
        except Exception as e:
            raise Exception(f"处理PDF文件时出错: {str(e)}")
            
        return documents

    def _process_docx(self, file_path: str) -> List[Document]:
        """处理DOCX文件"""
        if docx is None:
            raise ImportError("缺少python-docx库，无法处理DOCX文件")
            
        documents = []
        try:
            doc = docx.Document(file_path)
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
                    
            text = '\n'.join(full_text)
            if text.strip():
                tokenized_text = tokenize(text)
                documents.append(Document(
                    page_content=tokenized_text,
                    metadata={
                        'file_path': file_path,
                        'file_type': 'docx'
                    }
                ))
        except Exception as e:
            raise Exception(f"处理DOCX文件时出错: {str(e)}")
            
        return documents

    def _process_doc(self, file_path: str) -> List[Document]:
        """处理DOC文件（简化处理）"""
        # 简化处理，将DOC文件视为文本文件
        return self._process_txt(file_path)

    def _process_pptx(self, file_path: str) -> List[Document]:
        """处理PPTX文件"""
        if pptx is None:
            raise ImportError("缺少python-pptx库，无法处理PPTX文件")
            
        documents = []
        try:
            prs = pptx.Presentation(file_path)
            full_text = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    full_text.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_text))
                    
            text = '\n\n'.join(full_text)
            if text.strip():
                tokenized_text = tokenize(text)
                documents.append(Document(
                    page_content=tokenized_text,
                    metadata={
                        'file_path': file_path,
                        'file_type': 'pptx'
                    }
                ))
        except Exception as e:
            raise Exception(f"处理PPTX文件时出错: {str(e)}")
            
        return documents

    def _process_xlsx(self, file_path: str) -> List[Document]:
        """处理Excel文件"""
        if load_workbook is None:
            raise ImportError("缺少openpyxl库，无法处理Excel文件")
            
        documents = []
        try:
            wb = load_workbook(file_path, read_only=True)
            full_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = [f"Sheet: {sheet_name}"]
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        sheet_text.append(row_text)
                if len(sheet_text) > 1:  # 除了标题外还有内容
                    full_text.append("\n".join(sheet_text))
                    
            text = '\n\n'.join(full_text)
            if text.strip():
                tokenized_text = tokenize(text)
                documents.append(Document(
                    page_content=tokenized_text,
                    metadata={
                        'file_path': file_path,
                        'file_type': 'excel'
                    }
                ))
        except Exception as e:
            raise Exception(f"处理Excel文件时出错: {str(e)}")
            
        return documents

    def _process_html(self, file_path: str) -> List[Document]:
        """处理HTML文件"""
        if BeautifulSoup is None:
            raise ImportError("缺少beautifulsoup4库，无法处理HTML文件")
            
        documents = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                
            soup = BeautifulSoup(content, 'html.parser')
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.decompose()
                
            text = soup.get_text()
            # 清理文本
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            if text.strip():
                tokenized_text = tokenize(text)
                documents.append(Document(
                    page_content=tokenized_text,
                    metadata={
                        'file_path': file_path,
                        'file_type': 'html'
                    }
                ))
        except Exception as e:
            raise Exception(f"处理HTML文件时出错: {str(e)}")
            
        return documents

    def _process_txt(self, file_path: str) -> List[Document]:
        """处理文本文件"""
        documents = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                
            if text.strip():
                tokenized_text = tokenize(text)
                documents.append(Document(
                    page_content=tokenized_text,
                    metadata={
                        'file_path': file_path,
                        'file_type': 'text'
                    }
                ))
        except Exception as e:
            raise Exception(f"处理文本文件时出错: {str(e)}")
            
        return documents

    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式列表"""
        return list(self.supported_formats.keys())


def process_and_tokenize_file(file_path: str) -> List[Document]:
    """
    处理并分词单个文件的便捷函数
    
    Args:
        file_path: 文件路径
        
    Returns:
        List[Document]: 处理后的文档列表
    """
    processor = FileProcessor()
    return processor.process_file(file_path)


def batch_process_files(file_paths: List[str]) -> List[Document]:
    """
    批量处理多个文件
    
    Args:
        file_paths: 文件路径列表
        
    Returns:
        List[Document]: 所有文件处理后的文档列表
    """
    processor = FileProcessor()
    all_documents = []
    
    for file_path in file_paths:
        try:
            documents = processor.process_file(file_path)
            all_documents.extend(documents)
        except Exception as e:
            print(f"处理文件 {file_path} 时出错: {str(e)}")
            continue
            
    return all_documents


# 使用示例
if __name__ == "__main__":
    # 示例：处理单个文件
    processor = FileProcessor()
    
    # 显示支持的格式
    print("支持的文件格式:", processor.get_supported_formats())
    
    # 注意：实际使用时需要提供真实的文件路径
    # documents = processor.process_file("path/to/your/file.pdf")
    # for doc in documents:
    #     print(f"文件: {doc.metadata['file_path']}")
    #     print(f"内容: {doc.page_content[:100]}...")
    #     print("-" * 50)