#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
import sys
import os

# 假设你的项目目录是 /path/to/my_project
sys.path.append(os.path.abspath('/add/morefile'))

import copy
import re
from io import BytesIO
import zlib
from PIL import Image

# 尝试导入aspose.slides，如果失败则使用备用方案
try:
    import aspose.slides as slides
    import aspose.pydrawing as drawing
    ASPOSE_AVAILABLE = True
except ImportError:
    ASPOSE_AVAILABLE = False
    print("Warning: aspose.slides not found. Using python-pptx as fallback for PPT processing.")

from add.morefile.rag.nlp import tokenize, is_english
from add.morefile.rag.nlp import rag_tokenizer
from add.morefile.deepdoc.parser import PdfParser, PptParser, PlainParser

# 尝试导入python-pptx作为备用方案
try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not found. PPT processing may be limited.")

from PyPDF2 import PdfReader as pdf2_read


class Ppt(PptParser):
    def __call__(self, fnm, from_page, to_page, callback=None):
        txts = super().__call__(fnm, from_page, to_page)

        callback(0.5, "Text extraction finished.")
        
        imgs = []
        if ASPOSE_AVAILABLE:
            # 使用aspose.slides处理PPT
            import aspose.slides as slides
            import aspose.pydrawing as drawing
            with open(fnm, 'rb') as f:
                ppt_content = f.read()
            with slides.Presentation(BytesIO(ppt_content)) as presentation:
                for i, slide in enumerate(presentation.slides[from_page: to_page]):
                    buffered = BytesIO()
                    slide.get_thumbnail(
                        0.5, 0.5).save(
                        buffered, drawing.imaging.ImageFormat.jpeg)
                    imgs.append(Image.open(buffered))
        elif PPTX_AVAILABLE:
            # 使用python-pptx作为备用方案
            try:
                prs = pptx.Presentation(fnm)
                for i, slide in enumerate(prs.slides[from_page: to_page]):
                    # 创建一个简单的占位符图像
                    img = Image.new('RGB', (800, 600), color='white')
                    imgs.append(img)
            except Exception as e:
                print(f"Warning: Failed to process PPT with python-pptx: {e}")
                # 创建空白图像作为后备
                for i in range(len(txts)):
                    img = Image.new('RGB', (800, 600), color='white')
                    imgs.append(img)
        else:
            # 如果都没有，创建空白图像
            for i in range(len(txts)):
                img = Image.new('RGB', (800, 600), color='white')
                imgs.append(img)
                
        # 确保imgs和txts长度一致
        if len(imgs) != len(txts):
            # 如果长度不一致，用空白图像填充
            while len(imgs) < len(txts):
                img = Image.new('RGB', (800, 600), color='white')
                imgs.append(img)
            # 如果图像过多，截断
            if len(imgs) > len(txts):
                imgs = imgs[:len(txts)]
                
        callback(0.9, "Image extraction finished")
        self.is_english = is_english(txts)
        return list(zip(txts, imgs))


class Pdf(PdfParser):
    def __init__(self):
        super().__init__()

    def __garbage(self, txt):
        txt = txt.lower().strip()
        if re.match(r"[0-9\.,%/-]+$", txt):
            return True
        if len(txt) < 3:
            return True
        return False

    def __call__(self, filename, binary=None, from_page=0,
                 to_page=100000, zoomin=3, callback=None):
        callback(msg="OCR is running...")
        self.__images__(filename if not binary else binary,
                        zoomin, from_page, to_page, callback)
        callback(0.8, "Page {}~{}: OCR finished".format(
            from_page, min(to_page, self.total_page)))
        assert len(self.boxes) == len(self.page_images), "{} vs. {}".format(
            len(self.boxes), len(self.page_images))
        res = []
        for i in range(len(self.boxes)):
            lines = "\n".join([b["text"] for b in self.boxes[i]
                              if not self.__garbage(b["text"])])
            res.append((lines, self.page_images[i]))
        callback(0.9, "Page {}~{}: Parsing finished".format(
            from_page, min(to_page, self.total_page)))
        return res


class PlainPdf(PlainParser):
    def __call__(self, filename, binary=None, from_page=0,
                 to_page=100000, callback=None, **kwargs):
        self.pdf = pdf2_read(filename if not binary else BytesIO(binary))
        page_txt = []
        for page in self.pdf.pages[from_page: to_page]:
            page_txt.append(page.extract_text())
        callback(0.9, "Parsing finished")
        return [(txt, None) for txt in page_txt]


def chunk(filename, binary=None, from_page=0, to_page=100000,
          lang="Chinese", callback=None, **kwargs):
    """
    The supported file formats are pdf, pptx.
    Every page will be treated as a chunk. And the thumbnail of every page will be stored.
    PPT file will be parsed by using this method automatically, setting-up for every PPT file is not necessary.
    """
    eng = lang.lower() == "english"
    doc = {
        "docnm_kwd": filename,
        "title_tks": rag_tokenizer.tokenize(re.sub(r"\.[a-zA-Z]+$", "", filename))
    }
    doc["title_sm_tks"] = rag_tokenizer.fine_grained_tokenize(doc["title_tks"])
    res = []
    if re.search(r"\.pptx?$", filename, re.IGNORECASE):
        ppt_parser = Ppt()
        for pn, (txt, img) in enumerate(ppt_parser(
                filename if not binary else binary, from_page, 1000000, callback)):
            d = copy.deepcopy(doc)
            pn += from_page
            d["image"] = img
            d["page_num_int"] = [pn + 1]
            d["top_int"] = [0]
            d["position_int"] = [(pn + 1, 0, img.size[0], 0, img.size[1])]
            tokenize(d, txt, eng)
            #r += ''
            res.append(d)
        return res
    elif re.search(r"\.pdf$", filename, re.IGNORECASE):
        pdf_parser = Pdf() if kwargs.get(
            "parser_config", {}).get(
            "layout_recognize", True) else PlainPdf()
        for pn, (txt, img) in enumerate(pdf_parser(filename, binary,
                                                   from_page=from_page, to_page=to_page, callback=callback)):
            d = copy.deepcopy(doc)
            pn += from_page
            if img:
                d["image"] = img
            d["page_num_int"] = [pn + 1]
            d["top_int"] = [0]
            d["position_int"] = [
                (pn + 1, 0, img.size[0] if img else 0, 0, img.size[1] if img else 0)]
            tokenize(d, txt, eng)
            res.append(d)
        return res

    raise NotImplementedError(
        "file type not supported yet(pptx, pdf supported)")


if __name__ == "__main__":
    import sys

    def dummy(a, b):
        pass
    r = chunk('a.pptx', callback=dummy)

    #chunk(sys.argv[1], callback=dummy)
